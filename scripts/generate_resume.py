from pathlib import Path
from shutil import copyfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "resume.pptx"
DIST = ROOT / "dist" / "resume.pptx"
PROFILE = ROOT / "src" / "assets" / "images" / "profile_v4.png"


BG = RGBColor(0x09, 0x0B, 0x10)
SURFACE = RGBColor(0x11, 0x15, 0x1C)
SURFACE2 = RGBColor(0x17, 0x1C, 0x24)
BORDER = RGBColor(0x23, 0x2A, 0x35)
TEXT = RGBColor(0xF3, 0xEF, 0xE7)
MUTED = RGBColor(0xA9, 0xB0, 0xB8)
DIM = RGBColor(0x6F, 0x78, 0x82)
EMERALD = RGBColor(0x6F, 0xAE, 0x8D)
AMBER = RGBColor(0xC8, 0xA9, 0x6B)
SLATE = RGBColor(0x7E, 0x97, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_box(slide, x, y, w, h, fill, line=BORDER, radius=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(radius, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_textbox(slide, x, y, w, h, text=""):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if text:
        tf.text = text
    return box


def add_run(paragraph, text, *, font_name="Aptos", size=18, color=TEXT, bold=False, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return run


def style_paragraph(paragraph, *, level=0, space_after=0, line_spacing=1.15, align=PP_ALIGN.LEFT):
    paragraph.level = level
    paragraph.space_after = Pt(space_after)
    paragraph.line_spacing = line_spacing
    paragraph.alignment = align


def add_chip(slide, x, y, w, h, text, fill, color, border=None):
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    chip.fill.solid()
    chip.fill.fore_color.rgb = fill
    chip.line.color.rgb = border or fill
    chip.line.width = Pt(1)
    tf = chip.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    style_paragraph(p, align=PP_ALIGN.CENTER)
    add_run(p, text, font_name="Aptos", size=10, color=color, bold=True)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return chip


def add_section_title(slide, x, y, text):
    box = add_textbox(slide, x, y, Inches(3), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    style_paragraph(p, space_after=0)
    add_run(p, text.upper(), font_name="Aptos", size=9, color=DIM, bold=True)
    return box


def add_bullet_block(slide, x, y, w, h, items, font_size=11.5, bullet_color=EMERALD):
    box = add_textbox(slide, x, y, w, h)
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = ""
        style_paragraph(p, space_after=6, line_spacing=1.15)
        add_run(p, "• ", font_name="Aptos", size=font_size, color=bullet_color, bold=True)
        add_run(p, item, font_name="Aptos", size=font_size, color=MUTED)
    return box


def add_experience_card(slide, x, y, w, h, date, role, company, bullets, accent):
    card = add_box(slide, x, y, w, h, SURFACE, BORDER)
    # accent line
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Inches(0.14), y + Inches(0.16), Inches(0.05), h - Inches(0.32))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()

    date_box = add_textbox(slide, x + Inches(0.3), y + Inches(0.14), w - Inches(0.45), Inches(0.2))
    p = date_box.text_frame.paragraphs[0]
    style_paragraph(p)
    add_run(p, date, font_name="Aptos", size=9, color=DIM, bold=True)

    role_box = add_textbox(slide, x + Inches(0.3), y + Inches(0.37), w - Inches(0.45), Inches(0.3))
    p = role_box.text_frame.paragraphs[0]
    style_paragraph(p)
    add_run(p, role, font_name="Aptos", size=14.5, color=TEXT, bold=True)
    add_run(p, f"  {company}", font_name="Aptos", size=13.5, color=accent)

    add_bullet_block(slide, x + Inches(0.3), y + Inches(0.71), w - Inches(0.45), h - Inches(0.82), bullets, font_size=9.5)
    return card


def add_skill_pills(slide, x, y, w, labels):
    cur_x, cur_y = x, y
    row_h = Inches(0.28)
    gap = Inches(0.08)
    max_x = x + w
    for label in labels:
        est_w = Inches(0.16 + len(label) * 0.055)
        if cur_x + est_w > max_x:
            cur_x = x
            cur_y += row_h + gap
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, cur_x, cur_y, est_w, row_h)
        pill.fill.solid()
        pill.fill.fore_color.rgb = SURFACE2
        pill.line.color.rgb = BORDER
        pill.line.width = Pt(1)
        tf = pill.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        style_paragraph(p, align=PP_ALIGN.CENTER)
        add_run(p, label, font_name="Aptos", size=10, color=MUTED, bold=True)
        cur_x += est_w + gap


def add_two_col_keywords(slide, x, y, w, left_items, right_items):
    left = add_textbox(slide, x, y, w / 2 - Inches(0.08), Inches(1.05))
    right = add_textbox(slide, x + w / 2 + Inches(0.04), y, w / 2 - Inches(0.08), Inches(1.05))
    for box, items in [(left, left_items), (right, right_items)]:
        tf = box.text_frame
        tf.clear()
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            style_paragraph(p, space_after=4)
            add_run(p, item, font_name="Aptos", size=10.5, color=MUTED, bold=True)


def add_cert_list(slide, x, y, w, items):
    box = add_textbox(slide, x, y, w, Inches(1.55))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        style_paragraph(p, space_after=4, line_spacing=1.05)
        add_run(p, item[0], font_name="Aptos", size=10.5, color=TEXT, bold=True)
        add_run(p, f" — {item[1]}", font_name="Aptos", size=10, color=MUTED)
    return box


def build_resume():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG

    # outer frame
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.12), Inches(0.12), Inches(13.09), Inches(7.26))
    frame.fill.background()
    frame.line.color.rgb = BORDER
    frame.line.width = Pt(1)

    # header
    add_chip(slide, Inches(0.35), Inches(0.28), Inches(0.52), Inches(0.27), "K3", SURFACE, AMBER, AMBER)

    name_box = add_textbox(slide, Inches(0.92), Inches(0.23), Inches(4.2), Inches(0.34))
    p = name_box.text_frame.paragraphs[0]
    style_paragraph(p)
    add_run(p, "Kim Edrian Binasoy", font_name="Georgia", size=19, color=TEXT, bold=False)
    add_run(p, "   [ relearn.ing ]", font_name="Aptos", size=10.5, color=EMERALD)

    right_contact = add_textbox(slide, Inches(8.1), Inches(0.22), Inches(4.8), Inches(0.36))
    p = right_contact.text_frame.paragraphs[0]
    style_paragraph(p, align=PP_ALIGN.RIGHT)
    add_run(p, "hello@binasoy.kim  •  Toronto, Canada  •  github.com/kedbin", font_name="Aptos", size=10.5, color=MUTED)

    divider = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.32), Inches(0.68), Inches(12.7), Inches(0.01))
    divider.fill.solid()
    divider.fill.fore_color.rgb = BORDER
    divider.line.fill.background()

    # left column
    left_x, left_y, left_w, left_h = Inches(0.38), Inches(0.86), Inches(3.2), Inches(5.8)
    add_box(slide, left_x, left_y, left_w, left_h, SURFACE)

    # glow rings behind portrait
    for offs, color in [(0.0, EMERALD), (0.18, SLATE), (0.34, AMBER)]:
        ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left_x + Inches(0.42) - Inches(offs), left_y + Inches(0.82) - Inches(offs), Inches(2.18 + offs * 2), Inches(2.18 + offs * 2))
        ring.fill.background()
        ring.line.color.rgb = color
        ring.line.width = Pt(1)
        ring.line.transparency = 0.7

    slide.shapes.add_picture(str(PROFILE), left_x + Inches(0.55), left_y + Inches(0.98), width=Inches(2.0), height=Inches(2.98))

    add_chip(slide, left_x + Inches(0.28), left_y + Inches(0.24), Inches(1.05), Inches(0.28), "OPEN TO WORK", RGBColor(0x12, 0x23, 0x1B), EMERALD, EMERALD)

    role_box = add_textbox(slide, left_x + Inches(0.28), left_y + Inches(4.18), Inches(2.64), Inches(0.56))
    p = role_box.text_frame.paragraphs[0]
    style_paragraph(p, line_spacing=1.0)
    add_run(p, "Cloud Engineer", font_name="Georgia", size=23, color=TEXT, bold=False)
    p = role_box.text_frame.add_paragraph()
    style_paragraph(p)
    add_run(p, "Automation • AI Systems • Developer", font_name="Aptos", size=11.5, color=EMERALD, bold=True)

    summary_box = add_textbox(slide, left_x + Inches(0.28), left_y + Inches(4.9), Inches(2.55), Inches(0.9))
    p = summary_box.text_frame.paragraphs[0]
    style_paragraph(p, line_spacing=1.18)
    add_run(p, "Builds cloud-native systems, automation workflows, and AI-assisted tools that reduce friction, improve onboarding, and scale delivery.", font_name="Aptos", size=10.5, color=MUTED)

    # right main column
    rx = Inches(3.9)
    headline = add_textbox(slide, rx, Inches(0.92), Inches(5.45), Inches(0.82))
    p = headline.text_frame.paragraphs[0]
    style_paragraph(p, line_spacing=0.95)
    add_run(p, "I build ", font_name="Georgia", size=26, color=TEXT)
    add_run(p, "systems", font_name="Georgia", size=26, color=EMERALD)
    p = headline.text_frame.add_paragraph()
    style_paragraph(p, line_spacing=0.95)
    add_run(p, "that automate, scale, and ", font_name="Georgia", size=26, color=TEXT)
    add_run(p, "learn.", font_name="Georgia", size=26, color=AMBER)

    sub = add_textbox(slide, rx, Inches(1.72), Inches(5.0), Inches(0.4))
    p = sub.text_frame.paragraphs[0]
    style_paragraph(p)
    add_run(p, "Selected experience, certifications, and keywords optimized for human + LLM screening.", font_name="Aptos", size=10.5, color=MUTED)

    add_section_title(slide, rx, Inches(2.12), "Selected Experience")
    add_experience_card(
        slide, rx, Inches(2.33), Inches(5.35), Inches(1.18),
        "NOV 2025 — PRESENT", "Founder", "Relearn.ing",
        [
            "Leading product direction for a digital learning platform centered on AI-era work and continuous improvement.",
            "Architecting cloud-native workflows for content delivery, automation, and user engagement.",
            "Designing ethical AI-assisted systems that help users operationalize learning and personal growth."
        ],
        EMERALD,
    )
    add_experience_card(
        slide, rx, Inches(3.59), Inches(5.35), Inches(1.18),
        "JUN 2023 — NOV 2025", "Quality Automation Engineer", "Transportation",
        [
            "Executed 800+ test cases and shipped an automation tool that reduced documentation time by 70%.",
            "Improved testing efficiency by 40% through cross-functional automation and process improvements.",
            "Built MS Teams bots and GenAI-powered knowledge transfer tooling for onboarding and operations."
        ],
        SLATE,
    )
    add_experience_card(
        slide, rx, Inches(4.85), Inches(5.35), Inches(0.98),
        "2021 — 2023", "Program Mgmt & App Dev", "Banking & IT",
        [
            "Managed onboarding/offboarding workflows for 200+ resources and improved operational efficiency by 50–75%.",
            "Built Power Automate, PowerApps, Teams tooling, and automated web workflows across internal operations."
        ],
        AMBER,
    )

    # right side rail
    sx = Inches(9.55)
    add_section_title(slide, sx, Inches(0.98), "Core Stack")
    add_two_col_keywords(
        slide, sx, Inches(1.18), Inches(3.15),
        ["Azure", "AWS", "Python", "TypeScript", "React", "Node.js"],
        ["Power Automate", "Cloud Operations", "QA Automation", "Runbooks", "Incident Triage", "GenAI Tooling"],
    )

    add_section_title(slide, sx, Inches(2.4), "Certifications")
    add_cert_list(slide, sx, Inches(2.64), Inches(3.1), [
        ("Google Professional Cloud DevOps Engineer", "Feb 2026 — Feb 2028"),
        ("Google Professional Security Operations Engineer", "Jan 2026 — Jan 2028"),
        ("Google Professional Cloud Architect", "Jan 2026 — Jan 2028"),
        ("Google Generative AI Leader", "Dec 2025 — Dec 2028"),
        ("Azure Developer / AI Fundamentals / Azure Fundamentals / AWS CCP", "Additional certifications"),
    ])

    add_section_title(slide, sx, Inches(5.58), "Highlights")
    highlights = add_textbox(slide, sx, Inches(5.84), Inches(3.15), Inches(0.86))
    tf = highlights.text_frame
    for idx, line in enumerate([
        "4+ years across transportation, banking, and IT.",
        "800+ executed test cases.",
        "70% documentation time reduction via automation.",
        "40% testing efficiency improvement.",
        "5-Day AI Agents Intensive Course with Google (ULSA).",
        "Portfolio: binasoy.kim • Publication: relearn.ing",
    ]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        style_paragraph(p, space_after=4)
        add_run(p, "• ", font_name="Aptos", size=10, color=EMERALD, bold=True)
        add_run(p, line, font_name="Aptos", size=9.6, color=MUTED)

    prs.save(OUT)
    DIST.parent.mkdir(parents=True, exist_ok=True)
    copyfile(OUT, DIST)


if __name__ == "__main__":
    build_resume()
